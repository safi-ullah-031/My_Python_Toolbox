import os
import gradio as gr
from pptx import Presentation
from PIL import Image
from io import BytesIO
from reportlab.pdfgen import canvas
from tqdm import tqdm

# Output Directory
OUTPUT_DIR = "converted_files"
os.makedirs(OUTPUT_DIR, exist_ok=True)

def convert_pptx(file_path, output_format):
    if not file_path:
        return "No file uploaded", None

    try:
        # Load PowerPoint file
        prs = Presentation(file_path)
        filenames = []
        
        # Progress bar setup
        progress_bar = tqdm(total=len(prs.slides), desc="Converting slides")

        if output_format == "Images":
            output_folder = os.path.join(OUTPUT_DIR, "images")
            os.makedirs(output_folder, exist_ok=True)
            
            for i, slide in enumerate(prs.slides):
                img_path = os.path.join(output_folder, f"slide_{i+1}.png")
                slide_image = convert_slide_to_image(slide)
                slide_image.save(img_path, "PNG")
                filenames.append(img_path)
                
                # Update progress bar
                progress_bar.update(1)

            progress_bar.close()
            return f"Conversion completed. {len(filenames)} images generated.", filenames

        elif output_format == "PDF":
            pdf_path = os.path.join(OUTPUT_DIR, "converted.pdf")
            create_pdf_from_pptx(prs, pdf_path)
            
            progress_bar.close()
            return f"Conversion completed. PDF saved at {pdf_path}", [pdf_path]

    except Exception as e:
        return f"Error: {str(e)}", None


def convert_slide_to_image(slide):
    """ Converts a PowerPoint slide to an image (dummy function). """
    width, height = 1280, 720  # Define resolution
    img = Image.new("RGB", (width, height), (255, 255, 255))
    
    draw = Image.new("RGBA", (width, height), (255, 255, 255, 0))
    img.paste(draw, (0, 0), draw)
    return img


def create_pdf_from_pptx(prs, pdf_path):
    """ Creates a PDF from PowerPoint slides. """
    c = canvas.Canvas(pdf_path)
    
    for i, slide in enumerate(prs.slides):
        c.drawString(100, 750, f"Slide {i+1}")  # Placeholder for content
        c.showPage()
    
    c.save()


# Gradio Interface
with gr.Blocks() as demo:
    gr.Markdown("## 📂 PPTX Converter (Images/PDF)")

    with gr.Row():
        file_input = gr.File(label="Upload PPTX File", type="filepath")
        output_format = gr.Radio(["Images", "PDF"], label="Select Output Format")
    
    with gr.Row():
        convert_button = gr.Button("Convert")
    
    progress_label = gr.Label(value="Awaiting conversion...")
    gallery_output = gr.Gallery(label="Converted Files", columns=3, height=300)
    
    convert_button.click(convert_pptx, inputs=[file_input, output_format], outputs=[progress_label, gallery_output])

demo.launch()
